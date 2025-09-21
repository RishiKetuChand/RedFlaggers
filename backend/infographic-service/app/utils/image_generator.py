import logging
import os
import tempfile
import uuid
import shutil
import subprocess
from pathlib import Path
from typing import Dict, Any, List
from datetime import datetime

from pptx import Presentation

from google.cloud import storage

logger = logging.getLogger(__name__)


class ImageGenerator:
    def __init__(self):
        self.storage_client = storage.Client()

    def _replace_in_runs(self, slide, replacements: dict):
        """Replace placeholders only inside existing runs to preserve formatting/alignment."""
        def replace_in_text_frame(tf):
            if tf is None:
                return
            for para in tf.paragraphs:
                for run in para.runs:
                    if not run.text:
                        continue
                    t = run.text
                    for k, v in replacements.items():
                        if k in t:
                            t = t.replace(k, v if v is not None else "")
                    if t != run.text:
                        run.text = t

        for shp in slide.shapes:
            if hasattr(shp, "text_frame"):
                replace_in_text_frame(shp.text_frame)
            if getattr(shp, "has_table", False):
                for row in shp.table.rows:
                    for cell in row.cells:
                        replace_in_text_frame(cell.text_frame)

    def modify_template(self, payload: dict) -> str:
        """Modify PPTX template with infographic results and return path to modified file."""
        base_dir = os.path.dirname(__file__)  
        template_path = os.path.join(base_dir, "template.pptx")
        prs = Presentation(template_path)

        startup_name = payload.get("startup_name", "") or ""
        results = payload.get("results", {}) or {}

        # Slide 0
        if len(prs.slides) > 0:
            self._replace_in_runs(prs.slides[0], {
                "PLACEHOLDER_STARTUP_NAME": startup_name,
            })

        # helper: tolerant field getter
        def gv(d: dict, *keys):
            for k in keys:
                if k in d:
                    return d.get(k, "") or ""
            return ""

        # Slide 1: product
        product = results.get("product")
        if product and len(prs.slides) > 1:
            self._replace_in_runs(prs.slides[1], {
                "PLACEHOLDER_PROBLEM":   gv(product, "problem"),
                "PLACEHOLDER_CP":        gv(product, "problem_category"),
                "PLACEHOLDER_CA":        gv(product, "current_alternatives"),
                "PLACEHOLDER_WN":        gv(product, "why_now"),
                "PLACEHOLDER_PD":        gv(product, "product_details"),
                "PLACEHOLDER_RF":        gv(product, "replacement_for"),
            })

        # Slide 2: financial_metric
        fin = results.get("financial_metric")
        if fin and len(prs.slides) > 2:
            self._replace_in_runs(prs.slides[2], {
                "PLACEHOLDER_OVERVIEW":  gv(fin, "overview"),
                "PLACEHOLDER_GR":        gv(fin, "growth_rate"),
                "PLACEHOLDER_CE":        gv(fin, "capital_efficiency"),
                "PLACEHOLDER_VALUATION": gv(fin, "valuation"),
                "PLACEHOLDER_ANALYSIS":  gv(fin, "analysis"),
                "PLACEHOLDER_PM":        gv(fin, "profitability_margin"),
            })
        # Save modified presentation
        filename = f"{startup_name}_{uuid.uuid4().hex[:8]}_infographic.pptx"
        temp_dir = tempfile.gettempdir()
        output_path = os.path.join(temp_dir, filename)
        prs.save(output_path)
        return output_path

    def convert_pptx_to_images(self, pptx_path: str, out_dir: str, startup_name: str, upload_id: str, fmt: str = "PNG", dpi: int = 220) -> List[str]:
        """
        PPTX -> PDF (LibreOffice), then PDF pages -> images (pdf2image).

        Args:
            pptx_path: path to .ppt or .pptx
            out_dir: directory to write images into (created if missing)
            startup_name: name of startup for filename
            upload_id: upload identifier for filename
            fmt: "PNG" or "JPEG"
            dpi: rasterization DPI for PDF->image path (Linux/macOS)
            
        Returns:
            List of image file paths
        """
        pptx_path = str(Path(pptx_path).resolve())
        out_dir = Path(out_dir)
        out_dir.mkdir(parents=True, exist_ok=True)

        soffice = shutil.which("soffice") or shutil.which("libreoffice")
        if not soffice:
            raise RuntimeError("LibreOffice not found. Install it (e.g., sudo apt install libreoffice).")

        try:
            from pdf2image import convert_from_path  # pip install pdf2image
        except ImportError as e:
            raise RuntimeError("pdf2image is required. Install with: pip install pdf2image") from e

        # Poppler's pdftoppm is required by pdf2image
        if not shutil.which("pdftoppm"):
            raise RuntimeError("Poppler not found. Install it (e.g., sudo apt install poppler-utils).")

        # 1) Convert PPTX -> PDF via LibreOffice (headless)
        with tempfile.TemporaryDirectory() as tmpdir:
            tmpdir = Path(tmpdir)
            cmd = [
                soffice, "--headless", "--invisible", "--nocrashreport", "--nodefault",
                "--view", "--nolockcheck", "--nologo",
                "--convert-to", "pdf", "--outdir", str(tmpdir),
                pptx_path
            ]
            subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)

            # Find produced PDF (LibreOffice keeps the base name, changes extension to .pdf)
            produced = Path(pptx_path).with_suffix(".pdf").name
            pdf_path = tmpdir / produced
            if not pdf_path.exists():
                # Fallback: pick any PDF produced
                pdfs = list(tmpdir.glob("*.pdf"))
                if not pdfs:
                    raise RuntimeError("LibreOffice did not produce a PDF. Is the file valid?")
                pdf_path = pdfs[0]

            # 2) Convert each PDF page -> image
            images = convert_from_path(str(pdf_path), dpi=dpi)
            
            # Save images with proper naming
            image_paths = []
            for idx, img in enumerate(images, start=1):
                out_name = f"{startup_name}_{upload_id}_image_{idx}.{fmt.lower()}"
                image_path = out_dir / out_name
                img.save(image_path, fmt.upper())
                image_paths.append(str(image_path))
                
            return image_paths


    def upload_to_gcs(self, file_path: str, bucket_name: str, blob_name: str = None) -> str:
        """Upload file to Google Cloud Storage and return public URL"""
        if not blob_name:
            blob_name = f"results/{os.path.basename(file_path)}"
        
        try:
            # Get bucket
            bucket = self.storage_client.bucket(bucket_name)
            blob = bucket.blob(blob_name)
            
            # Determine content type based on file extension
            file_ext = os.path.splitext(file_path)[1].lower()
            content_type = 'image/png' if file_ext == '.png' else 'application/pdf'
            
            # Upload file
            with open(file_path, 'rb') as file_data:
                blob.upload_from_file(file_data, content_type=content_type)
            
            # Make blob publicly accessible
            blob.make_public()
            
            # Get public URL
            public_url = blob.public_url
            
            logger.info(f"File uploaded successfully to GCS: {public_url}")
            return public_url
            
        except Exception as e:
            logger.error(f"Failed to upload file to GCS: {str(e)}")
            raise e
        finally:
            # Clean up local file
            try:
                if os.path.exists(file_path):
                    #TODO remove it later
                    #os.remove(file_path)
                    logger.info(f"Cleaned up local file: {file_path}")
            except Exception as cleanup_error:
                logger.warning(f"Failed to cleanup local file: {str(cleanup_error)}")

    def process_infographic_to_images_url(self, infographic_data: Dict[str, Any], bucket_name: str, upload_id: str) -> Dict[str, List[str]]:
        """Complete workflow: infographic -> modify template -> convert to images -> upload -> return URLs"""
        try:
            startup_name = infographic_data.get('startup_name', 'unknown')
            
            # Create temporary directory for processing
            with tempfile.TemporaryDirectory() as temp_dir:
                # Modify template with infographic data
                pptx_path = self.modify_template(infographic_data)
                
                # Convert PPTX to images
                image_paths = self.convert_pptx_to_images(
                    pptx_path=pptx_path,
                    out_dir=temp_dir,
                    startup_name=startup_name,
                    upload_id=upload_id
                )
                
                # Upload images to GCS and collect URLs
                image_urls = []
                for i, image_path in enumerate(image_paths, 1):
                    filename = f"{startup_name}_{upload_id}_image_{i}.png"
                    blob_name = f"analysis_reports/{filename}"
                    
                    # Upload to GCS
                    public_url = self.upload_to_gcs(image_path, bucket_name, blob_name)
                    image_urls.append(public_url)
                
                # Clean up temporary PPTX file
                if os.path.exists(pptx_path):
                    os.remove(pptx_path)
                
                return {"image_urls": image_urls}
            
        except Exception as e:
            logger.error(f"Failed to process infographic to images: {str(e)}")
            raise e

